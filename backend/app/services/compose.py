"""AI 拼 PPT 服务(SE-06):大纲 → 逐项混合检索 → 取 top1 素材页 → 图片拼装新 PPT。

产出形态:每页一张满幅 preview PNG(图片版草稿,视觉/版式保真,文字不可编辑;
需要编辑的页可走单页 PPTX 导出替换)。依赖 python-pptx。
"""
from __future__ import annotations

import io
import logging
from dataclasses import dataclass, field

from sqlalchemy.orm import Session

from app.core.storage import get_storage
from app.models import Presentation, PresentationVersion, Slide, User
from app.services.hybrid_search import hybrid_search

logger = logging.getLogger(__name__)

# 16:9 页面尺寸(EMU)。python-pptx 默认 10x7.5in(4:3),改为 13.33x7.5in(16:9)。
_SLIDE_W = 12192000
_SLIDE_H = 6858000


@dataclass
class ComposeMatch:
    section: str
    query: str
    slide_id: str | None = None
    presentation_title: str | None = None
    page_no: int | None = None
    score: float = 0.0
    hit_reasons: list[str] = field(default_factory=list)
    matched: bool = False


def _match_outline(db: Session, outline: list[dict], user: User) -> list[ComposeMatch]:
    """对每个大纲项做混合检索,取 top1。visibility 由 hybrid_search 内部过滤。"""
    matches: list[ComposeMatch] = []
    for item in outline:
        section = str(item.get("section") or "").strip() or "未命名"
        query = str(item.get("query") or "").strip()
        m = ComposeMatch(section=section, query=query)
        if query:
            hits = hybrid_search(db, query, user_id=user.id, superuser=user.is_superuser, topn=1)
            if hits:
                h = hits[0]
                # 取所属 presentation 标题
                pres = (
                    db.query(Presentation.title)
                    .join(PresentationVersion, PresentationVersion.presentation_id == Presentation.id)
                    .filter(PresentationVersion.id == h.slide.version_id)
                    .scalar()
                )
                m.slide_id = h.slide.id
                m.presentation_title = pres
                m.page_no = h.slide.page_no
                m.score = round(h.score, 4)
                m.hit_reasons = h.hit_reasons
                m.matched = True
        matches.append(m)
    return matches


def _assemble_pptx(db: Session, matches: list[ComposeMatch]) -> bytes:
    """把匹配的素材页 preview PNG 按序拼装成新 PPTX(每页满幅图 + 备注来源)。"""
    from pptx import Presentation as PptxPresentation
    from pptx.util import Emu

    prs = PptxPresentation()
    prs.slide_width = Emu(_SLIDE_W)
    prs.slide_height = Emu(_SLIDE_H)
    blank = prs.slide_layouts[6]  # blank layout

    storage = get_storage()
    slide_ids = [m.slide_id for m in matches if m.slide_id]
    slides = {s.id: s for s in db.query(Slide).filter(Slide.id.in_(slide_ids)).all()} if slide_ids else {}

    for m in matches:
        slide = prs.slides.add_slide(blank)
        s = slides.get(m.slide_id) if m.slide_id else None
        png_bytes = None
        if s and s.preview_object_key:
            try:
                png_bytes = storage.get_object(s.preview_object_key)
            except Exception as e:  # noqa: BLE001
                logger.warning("preview fetch failed for slide %s: %s", m.slide_id, e)
        if png_bytes:
            slide.shapes.add_picture(io.BytesIO(png_bytes), 0, 0, width=Emu(_SLIDE_W), height=Emu(_SLIDE_H))
        # 备注:大纲 section + 来源
        notes = slide.notes_slide.notes_text_frame
        src = f"{m.presentation_title} P{m.page_no}" if m.matched else "(未找到素材,占位页)"
        notes.text = f"[{m.section}] {src}"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def compose_pptx(db: Session, title: str, outline: list[dict], user: User) -> tuple[list[ComposeMatch], bytes]:
    """主流程:大纲 → 匹配 → 拼装。返回 (matches 明细, pptx bytes)。"""
    matches = _match_outline(db, outline, user)
    pptx_bytes = _assemble_pptx(db, matches)
    return matches, pptx_bytes
